"""Document parser — 3-tier: Docling → Claude Vision → pdfplumber.

Flow:
  PDF → Docling → Markdown (always return, pipeline handles tables)
    - Docling quality bad → Vision per-page (context mode) → fallback pdfplumber
  DOCX → Docling → Markdown, fallback python-docx
  XLSX → openpyxl (structured text)
  TXT/MD → read directly
"""
from __future__ import annotations

import base64
import hashlib
import io
import logging
import mimetypes
import re
from datetime import datetime, timezone
from pathlib import Path
from typing import Union

logger = logging.getLogger(__name__)

_docling_converter = None

MIN_TEXT_CHARS = 50

# --- Image extraction config ---
# Bỏ ảnh nhỏ hơn ngưỡng này → loại logo/icon noise trước khi tốn Haiku call.
MIN_IMAGE_DIMENSION = 100  # px
# Bỏ ảnh quá lớn (thường là full-page scan đã có tier text song song) — caption Haiku
# cũng có giới hạn payload. 5MB là buffer rộng cho ảnh chụp camera.
MAX_IMAGE_BYTES = 5 * 1024 * 1024
# Cap số ảnh / trang để tránh slide hoa văn trang trí làm scale Haiku call vô tội vạ.
MAX_IMAGES_PER_PAGE = 10
# Caption ngắn 1-2 câu → 256 token đủ rộng.
CAPTION_MAX_TOKENS = 256

IMAGE_CAPTION_PROMPT = (
    "Mô tả hình ảnh này thành 1-2 câu tiếng Việt ngắn gọn, tập trung vào "
    "NỘI DUNG THÔNG TIN (sơ đồ gì, biểu đồ về cái gì, ảnh chụp cái gì, bảng nói về gì).\n"
    "KHÔNG mô tả style/format/màu sắc trang trí.\n"
    "Giữ chính xác tên riêng, số liệu, tiếng Việt có dấu.\n"
    "KHÔNG dùng Markdown.\n\n"
    "NGỮ CẢNH (text quanh ảnh trong tài liệu):\n"
    "---\n{context}\n---\n\n"
    "Chỉ trả về câu mô tả, không thêm giải thích."
)

VISION_TO_MARKDOWN_PROMPT = (
    "Hãy chuyển nội dung trang tài liệu này thành Markdown.\n"
    "Yêu cầu:\n"
    "- Nếu có bảng: chuyển thành Markdown table (| cột 1 | cột 2 |)\n"
    "- Nếu có tiêu đề: dùng ## heading\n"
    "- Nếu có danh sách: dùng - bullet\n"
    "- QUAN TRỌNG: Giữ nguyên chính xác tiếng Việt có dấu.\n"
    "  Ví dụ: 'Chống cháy' KHÔNG PHẢI 'Chông cháy', 'Cửa' KHÔNG PHẢI 'Của'\n"
    "- Chỉ trả về Markdown, không giải thích thêm."
)

VISION_TO_CONTEXT_PROMPT = (
    "Đọc trang tài liệu này và viết lại NỘI DUNG thành plain text có cấu trúc.\n"
    "Yêu cầu:\n"
    "- Dòng đầu: tiêu đề/mô tả ngắn về trang này.\n"
    "- Nếu có bảng: mỗi dòng dữ liệu viết thành 1 dòng text riêng.\n"
    "  VD: 'Video 1: Giới thiệu BKVN. Tình trạng: Lên kịch bản.'\n"
    "- Nếu ô có màu mang ý nghĩa (xanh=hoàn thành, cam=đang làm, đỏ=quan trọng):\n"
    "  → Ghi gắn với dữ liệu trên cùng dòng.\n"
    "  VD: 'Tất niên, du xuân (T1-T2): đã hoàn thành.'\n"
    "- Dữ liệu lặp: ghi 'Tháng 3, 4: lặp lại nội dung Tháng 2.'\n"
    "- Bỏ qua cột/ô trống.\n"
    "- Giữ chính xác tên riêng, số liệu, tiếng Việt có dấu.\n"
    "  VD: 'Chống cháy' KHÔNG PHẢI 'Chông cháy'\n"
    "- KHÔNG dùng Markdown (**, ##, `, |, ---).\n"
    "- KHÔNG mô tả format, font, kiểu bảng.\n"
    "- Chỉ trả về plain text."
)

# Common OCR/Vision typos
_TYPO_FIXES: list[tuple[str, str]] = [
    ("Chông Cháy", "Chống Cháy"), ("Chông cháy", "Chống cháy"),
    ("chông cháy", "chống cháy"), ("Chông Chảy", "Chống Cháy"),
    ("Chông chảy", "Chống cháy"), ("chông chảy", "chống cháy"),
    ("Của chống", "Cửa chống"), ("Của Chống", "Cửa Chống"),
    ("của chống", "cửa chống"), ("Của chông", "Cửa chống"),
    ("Của Chông", "Cửa Chống"), ("của chông", "cửa chống"),
    ("Dã quay", "Đã quay"), ("Dã dạng", "Đã đăng"),
    ("Dạng edit", "Đang edit"), ("Tính trạng", "Tình trạng"),
    ("KỀ HOẠCH", "KẾ HOẠCH"), ("Kề hoạch", "Kế hoạch"),
    ("kề hoạch", "kế hoạch"), ("ban nhạt định", "bạn nhất định"),
    ("Ban nhạt định", "Bạn nhất định"), ("phải ban nhạt", "phải bạn nhất"),
]


def _sha256(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()[:16]


def _fix_common_typos(text: str) -> str:
    for wrong, correct in _TYPO_FIXES:
        text = text.replace(wrong, correct)
    return text


def _clean_markdown(md: str) -> str:
    md = re.sub(r"<!--.*?-->", "", md, flags=re.DOTALL)
    md = re.sub(r"-{3,}\s*page\s*\d*\s*-{3,}", "", md, flags=re.IGNORECASE)
    # Remove Docling image placeholders (![image](...), ![Image 1](...), etc.)
    md = re.sub(r"!\[(?:image|Image\s*#?\d*|Figure\s*#?\d*|Picture\s*#?\d*|Photo\s*#?\d*)\]\([^)]*\)", "", md, flags=re.IGNORECASE)
    # Remove bare bracketed placeholders like [Image #1], [Figure 2]
    md = re.sub(r"\[(?:Image|Figure|Picture|Photo)\s*#?\d+\](?!\()", "", md, flags=re.IGNORECASE)
    md = re.sub(r"\n{3,}", "\n\n", md)
    lines = [line.rstrip() for line in md.split("\n")]
    return "\n".join(lines).strip()


def _fix_joined_words(text: str) -> str:
    text = re.sub(
        r"([a-zàáảãạăắằẳẵặâấầẩẫậèéẻẽẹêếềểễệìíỉĩịòóỏõọôốồổỗộơớờởỡợùúủũụưứừửữựỳýỷỹỵđ])"
        r"([A-ZÀÁẢÃẠĂẮẰẲẴẶÂẤẦẨẪẬÈÉẺẼẸÊẾỀỂỄỆÌÍỈĨỊÒÓỎÕỌÔỐỒỔỖỘƠỚỜỞỠỢÙÚỦŨỤƯỨỪỬỮỰỲÝỶỸỴĐ])",
        r"\1 \2", text,
    )
    return text


def _count_meaningful_chars(text: str) -> int:
    return len(re.sub(r"[|#\-=_*`~>\s]", "", text))


def _check_docling_quality(md: str) -> tuple[bool, str]:
    if not md or not md.strip():
        return False, "empty output"
    lines = md.strip().split("\n")
    non_empty = [l for l in lines if l.strip()]
    if len(non_empty) <= 2 and md.count("|") > 20:
        return False, f"broken table: {md.count('|')} pipes on {len(non_empty)} lines"
    cells = md[:200].split("|")
    if len(cells) >= 3:
        first_cell = cells[1].strip() if len(cells) > 1 else ""
        if first_cell and len(first_cell) > 5 and md.count(first_cell) >= 5:
            return False, f"repeated header x{md.count(first_cell)}"
    if len(md) > 500 and len(non_empty) < 3:
        return False, f"no line breaks: {len(md)} chars in {len(non_empty)} lines"
    if len(md) > 500:
        paragraphs = [p.strip() for p in re.split(r"\n{2,}", md) if len(p.strip()) > 100]
        if len(paragraphs) >= 4:
            unique = set(paragraphs)
            dup_ratio = 1 - (len(unique) / len(paragraphs))
            if dup_ratio > 0.4:
                return False, f"duplicate blocks: {dup_ratio:.0%}"
    return True, "ok"


# --- Tier 1: Docling ---

def _get_docling_converter():
    global _docling_converter
    if _docling_converter is not None:
        return _docling_converter
    try:
        from docling.document_converter import DocumentConverter
        _docling_converter = DocumentConverter()
        logger.info("Docling converter loaded")
    except Exception as exc:
        logger.warning("Docling not available: %s", exc)
    return _docling_converter


def _docling_split_pages(doc) -> list[dict]:
    """Extract per-page content from Docling DoclingDocument. Returns [] on failure."""
    pages_content: dict[int, list[str]] = {}
    try:
        iterate = getattr(doc, "iterate_items", None)
        if iterate is None:
            return []
        for item, _level in iterate():
            if "Picture" in type(item).__name__ or "Image" in type(item).__name__:
                continue
            provs = getattr(item, "prov", None)
            if not provs:
                continue
            page_no = getattr(provs[0], "page_no", None)
            if not isinstance(page_no, int):
                continue
            text = ""
            if hasattr(item, "export_to_markdown"):
                try:
                    text = item.export_to_markdown() or ""
                except Exception:
                    text = getattr(item, "text", "") or ""
            else:
                text = getattr(item, "text", "") or ""
            text = text.strip()
            if text:
                pages_content.setdefault(page_no, []).append(text)
    except Exception:
        logger.warning("Docling page iteration failed, falling back to single-page")
        return []

    result: list[dict] = []
    for page_no in sorted(pages_content.keys()):
        text = "\n\n".join(pages_content[page_no])
        text = _clean_markdown(text)
        text = _fix_joined_words(text)
        if _count_meaningful_chars(text) >= MIN_TEXT_CHARS:
            result.append({"page": page_no, "text": text})
    return result


def _parse_with_docling(path: Path) -> list[dict] | None:
    """Generic Docling parser for any supported format (PDF, DOCX, PPTX, ...).

    Returns per-page list[{"page": N, "text": md}] or None on failure/empty.
    Callers apply their own additional quality checks if needed.
    """
    converter = _get_docling_converter()
    if converter is None:
        return None
    try:
        result = converter.convert(source=str(path))
        doc = result.document
        full_md = _clean_markdown(_fix_joined_words(doc.export_to_markdown()))
        if not full_md or _count_meaningful_chars(full_md) < MIN_TEXT_CHARS:
            return None
        num_pages = len(getattr(doc, "pages", {}))
        if num_pages > 1:
            pages = _docling_split_pages(doc)
            if pages:
                logger.info("Docling %s: %d pages", path.name, len(pages))
                return pages
        logger.info("Docling %s: 1 page, %d chars", path.name, len(full_md))
        return [{"page": 1, "text": full_md}]
    except Exception:
        logger.exception("Docling failed for %s", path.name)
        return None


def _docling_parse(path: Path) -> list[dict]:
    """PDF-specific Docling parser with quality check. Returns [] on failure/rejection."""
    pages = _parse_with_docling(path)
    if not pages:
        return []
    full_text = "\n\n".join(p["text"] for p in pages)
    is_good, reason = _check_docling_quality(full_text)
    if not is_good:
        logger.warning("Docling REJECTED for %s: %s", path.name, reason)
        return []
    return pages


# --- Tier 2: Claude Vision ---

def _get_anthropic_client():
    from app.config import ANTHROPIC_API_KEY
    if not ANTHROPIC_API_KEY:
        return None
    try:
        import anthropic
        return anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
    except Exception:
        return None


def _pdf_page_to_base64(path: Path, page_num: int) -> str | None:
    try:
        from pdf2image import convert_from_path
        images = convert_from_path(str(path), dpi=200, first_page=page_num, last_page=page_num)
        if not images:
            return None
        import io
        buf = io.BytesIO()
        images[0].save(buf, format="PNG")
        b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
        del images
        return b64
    except Exception:
        logger.exception("Failed to convert page %d to image", page_num)
        return None


def _vision_call(client, image_b64: str, model: str, prompt: str) -> str:
    try:
        response = client.messages.create(
            model=model, max_tokens=4096, temperature=0.1,
            messages=[{"role": "user", "content": [
                {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": image_b64}},
                {"type": "text", "text": prompt},
            ]}],
        )
        result = response.content[0].text.strip()
        return _fix_common_typos(result)
    except Exception:
        logger.exception("Claude Vision call failed")
        return ""


def _vision_parse_pdf(path: Path, to_context: bool = False) -> list[dict]:
    from app.config import CLAUDE_HAIKU_MODEL
    client = _get_anthropic_client()
    if client is None:
        return []
    num_pages = _count_pdf_pages(path)
    prompt = VISION_TO_CONTEXT_PROMPT if to_context else VISION_TO_MARKDOWN_PROMPT
    logger.info("Vision parsing %s: %d pages, context=%s", path.name, num_pages, to_context)
    pages: list[dict] = []
    for page_num in range(1, num_pages + 1):
        image_b64 = _pdf_page_to_base64(path, page_num)
        if not image_b64:
            pages.append({"page": page_num, "text": ""})
            continue
        text = _vision_call(client, image_b64, CLAUDE_HAIKU_MODEL, prompt)
        text = _clean_markdown(text)
        logger.info("Vision page %d: %d chars", page_num, len(text))
        pages.append({"page": page_num, "text": text})
    return pages


# --- Tier 3: pdfplumber ---

def _pdfplumber_parse(path: Path) -> list[dict]:
    try:
        import pdfplumber
    except ImportError:
        return []
    pages: list[dict] = []
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                try:
                    text = (page.extract_text() or "").strip()
                except Exception:
                    text = ""
                pages.append({"page": i, "text": text})
    except Exception:
        logger.exception("pdfplumber cannot open %s", path.name)
    return pages


def _count_pdf_pages(path: Path) -> int:
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            return len(pdf.pages)
    except Exception:
        return 1


# --- Image extraction + caption + S3 upload (PyMuPDF, 1-pass) ---

def _png_bytes_from_image(image_bytes: bytes) -> tuple[bytes, int, int] | None:
    """Convert ảnh raw bytes (jpg/png/...) sang PNG + check dimension.

    Trả (png_bytes, width, height) hoặc None nếu fail / dimension < min.
    """
    try:
        from PIL import Image
    except ImportError:
        logger.warning("Pillow chưa cài — bỏ qua image extraction")
        return None

    try:
        pil_img = Image.open(io.BytesIO(image_bytes))
        width, height = pil_img.size
    except Exception:
        logger.exception("PIL không decode được image")
        return None

    if width < MIN_IMAGE_DIMENSION or height < MIN_IMAGE_DIMENSION:
        return None

    if pil_img.mode in ("CMYK", "RGBA", "P", "LA"):
        pil_img = pil_img.convert("RGB")

    buf = io.BytesIO()
    try:
        pil_img.save(buf, format="PNG")
    except Exception:
        logger.exception("Không encode được PNG")
        return None
    return buf.getvalue(), width, height


def _haiku_caption(client, png_bytes: bytes, context_text: str) -> str:
    """Gọi Haiku Vision với ảnh + ngữ cảnh trang → caption tiếng Việt 1-2 câu."""
    from app.config import CLAUDE_HAIKU_MODEL
    img_b64 = base64.b64encode(png_bytes).decode("utf-8")
    ctx = context_text[:400] if context_text else "(không có ngữ cảnh)"
    prompt = IMAGE_CAPTION_PROMPT.format(context=ctx)
    try:
        response = client.messages.create(
            model=CLAUDE_HAIKU_MODEL,
            max_tokens=CAPTION_MAX_TOKENS,
            temperature=0.1,
            messages=[{"role": "user", "content": [
                {"type": "image", "source": {"type": "base64",
                                             "media_type": "image/png",
                                             "data": img_b64}},
                {"type": "text", "text": prompt},
            ]}],
        )
        caption = response.content[0].text.strip()
        return _fix_common_typos(caption)
    except Exception:
        logger.exception("Caption Haiku call failed")
        return ""


def _iter_pdf_images(pdf_path: Path):
    """Generator yield (page_num, ord_idx, image_bytes) cho mỗi ảnh trong PDF.

    PDF có concept page rõ ràng → page_num = thứ tự trang thật trong file.
    Cap MAX_IMAGES_PER_PAGE per trang để tránh slide hoa văn trang trí.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.warning("PyMuPDF chưa cài — bỏ qua image extraction")
        return

    try:
        doc = fitz.open(str(pdf_path))
    except Exception:
        logger.exception("PyMuPDF không mở được %s", pdf_path.name)
        return

    try:
        for page_idx in range(len(doc)):
            page = doc[page_idx]
            page_num = page_idx + 1
            try:
                img_list = page.get_images(full=True)
            except Exception:
                logger.exception("get_images failed on page %d", page_num)
                continue

            count = 0
            for ord_idx, img_info in enumerate(img_list):
                if count >= MAX_IMAGES_PER_PAGE:
                    logger.info("Page %d: hit MAX_IMAGES_PER_PAGE=%d, skip rest",
                                page_num, MAX_IMAGES_PER_PAGE)
                    break
                xref = img_info[0]
                try:
                    base = doc.extract_image(xref)
                except Exception:
                    logger.exception("extract_image failed xref=%s page=%d", xref, page_num)
                    continue
                image_bytes = base.get("image", b"")
                if image_bytes:
                    yield (page_num, ord_idx, image_bytes)
                    count += 1
    finally:
        doc.close()


def _iter_docx_images(docx_path: Path):
    """Generator yield (page_num, ord_idx, image_bytes) cho mỗi ảnh trong DOCX.

    DOCX không có page boundary rõ ràng (page do Word render runtime, không phải
    property của file). v1: gom tất cả ảnh vào page_num=1 → caption sẽ dùng
    context page=1 (text Docling parse cho page=1).

    Đọc qua python-docx `document.part.related_parts` — tất cả ảnh nhúng đều
    nằm trong relationships của phần document chính. Kèm cả header/footer images
    nếu có (bằng walk thêm header_part / footer_part).
    """
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx chưa cài — bỏ qua DOCX image extraction")
        return

    try:
        doc = Document(str(docx_path))
    except Exception:
        logger.exception("python-docx không mở được %s", docx_path.name)
        return

    ord_idx = 0
    yielded = 0
    for rid, part in doc.part.related_parts.items():
        if yielded >= MAX_IMAGES_PER_PAGE:
            logger.info("DOCX %s: hit MAX_IMAGES_PER_PAGE=%d, skip rest",
                        docx_path.name, MAX_IMAGES_PER_PAGE)
            break
        ctype = getattr(part, "content_type", "") or ""
        if not ctype.startswith("image/"):
            continue
        try:
            blob = part.blob
        except Exception:
            logger.exception("DOCX rel %s không đọc được blob", rid)
            continue
        if not blob:
            continue
        yield (1, ord_idx, blob)
        ord_idx += 1
        yielded += 1


def _iter_xlsx_images(xlsx_path: Path):
    """Generator yield (sheet_idx_1based, ord_idx, image_bytes) cho mỗi ảnh
    trong workbook XLSX.

    XLSX có concept sheet (tab) thay cho page → dùng sheet_idx (1-based) làm
    "page" để fit chung pipeline với PDF/DOCX. Sheet name lưu riêng ở
    `parse_xlsx` (qua field `sheet_name`).

    Dùng `openpyxl` API: `ws._images` — private nhưng stable từ 3.x. Mỗi
    Image object truy cập bytes qua `_data()` (callable). Cap
    MAX_IMAGES_PER_PAGE mỗi sheet để tránh dashboard có hàng chục icon.
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.warning("openpyxl chưa cài — bỏ qua XLSX image extraction")
        return

    try:
        # `read_only=False` bắt buộc — read_only mode KHÔNG load images.
        wb = load_workbook(str(xlsx_path), data_only=True, read_only=False)
    except Exception:
        logger.exception("openpyxl không mở được %s", xlsx_path.name)
        return

    try:
        for sheet_idx, sheet_name in enumerate(wb.sheetnames, start=1):
            ws = wb[sheet_name]
            images = getattr(ws, "_images", None) or []
            count = 0
            for ord_idx, image in enumerate(images):
                if count >= MAX_IMAGES_PER_PAGE:
                    logger.info("XLSX %s sheet %r: hit MAX_IMAGES_PER_PAGE=%d, skip rest",
                                xlsx_path.name, sheet_name, MAX_IMAGES_PER_PAGE)
                    break
                blob: bytes = b""
                # openpyxl Image: `_data` thường là callable trả bytes; vài
                # version expose `ref` (PIL Image) hoặc `.path` (BytesIO).
                try:
                    data_attr = getattr(image, "_data", None)
                    if callable(data_attr):
                        blob = data_attr() or b""
                    elif hasattr(image, "ref"):
                        from io import BytesIO
                        bio = BytesIO()
                        image.ref.save(bio, format="PNG")
                        blob = bio.getvalue()
                except Exception:
                    logger.exception("XLSX %s sheet %r image %d: không lấy được bytes",
                                     xlsx_path.name, sheet_name, ord_idx)
                    continue
                if not blob:
                    continue
                yield (sheet_idx, ord_idx, blob)
                count += 1
    finally:
        try:
            wb.close()
        except Exception:
            pass


def _process_image_iter(
    image_iter,
    doc_id: str,
    page_texts: dict[int, str],
    source_label: str = "doc",
) -> dict[int, list[dict]]:
    """Generic: nhận iterator yield (page, ord, bytes) → filter dimension/bytes
    → dedupe → caption Haiku per-page-context → upload S3 → trả by_page meta.

    Format-agnostic: PDF/DOCX/PPTX đều dùng chung qua iterator riêng. Chỉ khác
    nhau ở cách yield (page, ord, bytes) — phần xử lý ảnh sau đó dùng chung.

    `source_label` chỉ để log dễ đọc (vd "PDF test.pdf").
    """
    from app.core import s3_client
    if not s3_client.is_configured():
        logger.warning("S3 chưa config — bỏ qua image extraction")
        return {}

    anth_client = _get_anthropic_client()
    if anth_client is None:
        logger.warning("Anthropic unavailable — extract+upload nhưng caption rỗng")

    # Clean ảnh stale của doc cũ trên S3.
    s3_client.delete_doc_images(doc_id)

    seen_image_ids: set[str] = set()
    cache_url: dict[str, str] = {}
    cache_caption: dict[str, str] = {}
    by_page: dict[int, list[dict]] = {}
    page_count: dict[int, int] = {}  # log per-page count

    for page_num, ord_idx, image_bytes in image_iter:
        if len(image_bytes) > MAX_IMAGE_BYTES:
            continue

        image_id = hashlib.sha256(image_bytes).hexdigest()[:16]

        # Dedupe: ảnh đã xử lý (vd logo header lặp) → reuse url+caption, vẫn ghi
        # metadata cho page hiện tại để retrieval gắn đúng context.
        if image_id in seen_image_ids:
            by_page.setdefault(page_num, []).append({
                "image_id": image_id,
                "url": cache_url.get(image_id, ""),
                "caption": cache_caption.get(image_id, ""),
                "page": page_num,
                "ord": ord_idx,
                "width": 0,
                "height": 0,
            })
            continue

        # Convert sang PNG + check dimension.
        converted = _png_bytes_from_image(image_bytes)
        if converted is None:
            continue
        png_bytes, width, height = converted

        # Caption ngay (in-memory) — cần bytes để gọi Haiku.
        ctx = (page_texts.get(page_num) or "").strip()
        caption = ""
        if anth_client is not None:
            caption = _haiku_caption(anth_client, png_bytes, ctx)
            if caption:
                logger.info("Captioned image %s (page %d): %d chars",
                            image_id, page_num, len(caption))

        # Upload S3.
        url = s3_client.upload_image(png_bytes, doc_id, image_id)
        if not url:
            logger.warning("Skip image %s vì S3 upload fail", image_id)
            continue

        seen_image_ids.add(image_id)
        cache_url[image_id] = url
        cache_caption[image_id] = caption

        by_page.setdefault(page_num, []).append({
            "image_id": image_id,
            "url": url,
            "caption": caption,
            "page": page_num,
            "ord": ord_idx,
            "width": width,
            "height": height,
        })
        page_count[page_num] = page_count.get(page_num, 0) + 1

    for p, n in sorted(page_count.items()):
        logger.info("%s page %d: %d images uploaded", source_label, p, n)

    return by_page


def _extract_caption_upload_images(
    pdf_path: Path,
    doc_id: str,
    page_texts: dict[int, str],
) -> dict[int, list[dict]]:
    """PDF wrapper — giữ tên cũ cho backward compat. Internally gọi
    `_process_image_iter` với `_iter_pdf_images` extractor."""
    return _process_image_iter(
        _iter_pdf_images(pdf_path), doc_id, page_texts,
        source_label=f"PDF {pdf_path.name}",
    )


def _extract_caption_upload_images_docx(
    docx_path: Path,
    doc_id: str,
    page_texts: dict[int, str],
) -> dict[int, list[dict]]:
    """DOCX wrapper — gom ảnh vào page_num=1 (DOCX không có page property).
    Caption dùng context của page 1 (text Docling parse cho page đầu tiên)."""
    return _process_image_iter(
        _iter_docx_images(docx_path), doc_id, page_texts,
        source_label=f"DOCX {docx_path.name}",
    )


def _extract_caption_upload_images_xlsx(
    xlsx_path: Path,
    doc_id: str,
    page_texts: dict[int, str],
) -> dict[int, list[dict]]:
    """XLSX wrapper — "page" map sang sheet_idx (1-based). page_texts key
    cũng là sheet_idx, value là text snippet của sheet đó (header + vài row
    đầu — caller phải truncate để Haiku context không bị nhồi)."""
    return _process_image_iter(
        _iter_xlsx_images(xlsx_path), doc_id, page_texts,
        source_label=f"XLSX {xlsx_path.name}",
    )


# --- Main PDF parser ---

def _truncate_pdf(path: Path, max_pages: int) -> Path | None:
    """Trả về tmp file PDF chỉ giữ N trang đầu. None nếu không cắt nổi."""
    try:
        from pypdf import PdfReader, PdfWriter
    except ImportError:
        try:
            from PyPDF2 import PdfReader, PdfWriter  # type: ignore
        except ImportError:
            logger.warning("pypdf/PyPDF2 not installed — preview parses full PDF")
            return None
    try:
        import tempfile
        reader = PdfReader(str(path))
        if len(reader.pages) <= max_pages:
            return None
        writer = PdfWriter()
        for i in range(min(max_pages, len(reader.pages))):
            writer.add_page(reader.pages[i])
        tmp = Path(tempfile.mktemp(suffix=".pdf"))
        with tmp.open("wb") as f:
            writer.write(f)
        return tmp
    except Exception:
        logger.exception("Failed to truncate PDF for preview")
        return None


def _parse_pdf_text_only(path: Path) -> list[dict]:
    """Chạy 3-tier text parse, không extract ảnh. Dùng cho preview metadata
    (preview chỉ cần text cho LLM gen tiêu đề/tag — không lưu ảnh vào Qdrant).

    Trả về [{"page": int, "text": str, "images": []}] — schema tương thích với
    parse_pdf() full.
    """
    pages = _docling_parse(path)
    if pages:
        return [{**p, "images": []} for p in pages]
    vision_pages = _vision_parse_pdf(path, to_context=True)
    if any(p["text"].strip() for p in vision_pages):
        return [{**p, "images": []} for p in vision_pages]
    pages = _pdfplumber_parse(path) or [{"page": 1, "text": ""}]
    return [{**p, "images": []} for p in pages]


def parse_pdf(path: Path, doc_id: str = "", max_pages: int | None = None) -> list[dict]:
    """Parse PDF qua 3-tier text + extract ảnh độc lập bằng PyMuPDF.

    Returns: [{"page": int, "text": str, "images": list[dict]}].

    Khi `max_pages` set: preview mode — cắt PDF, KHÔNG extract ảnh (preview
    chỉ dùng cho metadata gen, không persist).
    Khi `max_pages` None (mặc định): full ingest — extract + caption + gắn ảnh
    đúng page (tận dụng `_docling_split_pages` trả per-page text).
    """
    # --- Preview path: bypass image extraction ---
    if max_pages and max_pages > 0:
        truncated = _truncate_pdf(path, max_pages)
        if truncated is not None:
            try:
                return _parse_pdf_text_only(truncated)
            finally:
                truncated.unlink(missing_ok=True)
        # Truncate fail (file ngắn hơn max_pages, hoặc pypdf không cài) → vẫn
        # parse text-only trên file gốc.
        return _parse_pdf_text_only(path)

    # --- Full mode: text trước, sau đó extract+caption+upload với context ---
    # Reorder so với spike v1: trước extract trước rồi caption sau, dẫn tới
    # caption không có context page khi Docling thắng (text 1 cục markdown).
    # Giờ: chạy text tier trước → có per-page text → caption với context đúng
    # → upload S3 ngay trong cùng pass.

    def _attach_images(pages: list[dict]) -> list[dict]:
        if not doc_id:
            for p in pages:
                p.setdefault("images", [])
            return pages
        page_texts = {p["page"]: p["text"] for p in pages}
        images_by_page = _extract_caption_upload_images(path, doc_id, page_texts)
        for p in pages:
            p["images"] = images_by_page.get(p["page"], [])
        return pages

    # Tier 1: Docling (per-page list nhờ refactor _docling_split_pages).
    pages = _docling_parse(path)
    if pages:
        return _attach_images(pages)

    # Tier 2: Vision per-page
    logger.info("PDF %s: Docling failed, trying Vision", path.name)
    vision_pages = _vision_parse_pdf(path, to_context=True)
    if any(p["text"].strip() for p in vision_pages):
        return _attach_images(vision_pages)

    # Tier 3: pdfplumber
    logger.warning("PDF %s: Vision failed, using pdfplumber", path.name)
    pages = _pdfplumber_parse(path) or [{"page": 1, "text": ""}]
    return _attach_images(pages)


# --- XLSX (kept from teammate's code) ---

def _xlsx_sheet_to_text(ws, max_rows: int | None = None) -> str:
    """Worksheet → plain text: 1 row = 1 block "Header: value\\nHeader2: value2".
    Tách helper để parse_xlsx full/preview chia sẻ logic format."""
    if max_rows and max_rows > 0:
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= max_rows:
                break
            rows.append(row)
    else:
        rows = list(ws.iter_rows(values_only=True))
    if not rows:
        return ""

    # Detect header: row đầu có ≥2 cell non-empty.
    header_idx = 0
    headers: list[str] = []
    for i, row in enumerate(rows):
        non_empty = [c for c in row if c is not None and str(c).strip()]
        if len(non_empty) >= 2:
            headers = [str(c).strip() if c else f"Col{j+1}" for j, c in enumerate(row)]
            header_idx = i
            break

    parts: list[str] = []
    if not headers:
        for row in rows:
            line = " | ".join(str(c).strip() for c in row if c is not None and str(c).strip())
            if line:
                parts.append(line)
        return "\n\n".join(parts)

    for row in rows[header_idx + 1:]:
        cells = [str(c).strip() if c is not None else "" for c in row]
        if not any(cells):
            continue
        row_parts = [f"{h}: {c}" for h, c in zip(headers, cells) if c]
        if row_parts:
            parts.append("\n".join(row_parts))
    return "\n\n".join(parts)


def parse_xlsx(
    path: Path,
    doc_id: str = "",
    max_rows: int | None = None,
) -> list[dict] | str:
    """Parse XLSX → text + extract ảnh nhúng (qua openpyxl).

    Full mode (`max_rows=None`): trả `[{"page": sheet_idx, "text": sheet_text,
    "sheet_name": str, "images": list}]` — schema đồng nhất với parse_pdf/docx.
    "page" = sheet index 1-based (sheet đầu = page 1) để fit chung pipeline.
    `sheet_name` field thêm để retrieval/UI hiển thị "Sheet: Doanh thu Q3"
    thay vì "page 2" — chỉ XLSX có field này, payload format không break PDF/DOCX.

    Preview mode (`max_rows>0`): trả `str` flat — backward compat cho metadata
    gen, KHÔNG extract ảnh.

    Image caption context: dùng snippet đầu mỗi sheet (~2000 chars) làm
    `page_texts` cho Haiku. Sheet 1000+ rows không nhồi hết vào prompt.
    """
    from openpyxl import load_workbook

    # Preview path: read_only nhanh hơn, không cần images.
    if max_rows and max_rows > 0:
        wb = load_workbook(str(path), data_only=True, read_only=True)
        try:
            parts: list[str] = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_text = _xlsx_sheet_to_text(ws, max_rows=max_rows)
                if not sheet_text:
                    continue
                if len(wb.sheetnames) > 1:
                    parts.append(f"# Sheet: {sheet_name}")
                parts.append(sheet_text)
            return "\n\n".join(parts)
        finally:
            try:
                wb.close()
            except Exception:
                pass

    # Full mode — read_only=False để openpyxl load `ws._images` cho extraction.
    wb = load_workbook(str(path), data_only=True, read_only=False)
    try:
        pages: list[dict] = []
        for sheet_idx, sheet_name in enumerate(wb.sheetnames, start=1):
            ws = wb[sheet_name]
            sheet_text = _xlsx_sheet_to_text(ws, max_rows=None)
            pages.append({
                "page": sheet_idx,
                "text": sheet_text,
                "sheet_name": sheet_name,
            })
    finally:
        try:
            wb.close()
        except Exception:
            pass

    if not doc_id:
        # Không có doc_id (vd preview qua kênh khác) → bỏ qua image extraction.
        for p in pages:
            p.setdefault("images", [])
        return pages

    # Truncate sheet_text → context ngắn cho Haiku caption (header + vài row).
    page_texts = {p["page"]: (p["text"][:2000] if p["text"] else "") for p in pages}
    images_by_page = _extract_caption_upload_images_xlsx(path, doc_id, page_texts)
    for p in pages:
        p["images"] = images_by_page.get(p["page"], [])
    return pages


# --- DOCX, TXT, MD ---

def parse_docx(
    path: Path,
    doc_id: str = "",
    max_paragraphs: int | None = None,
) -> list[dict] | str:
    """Parse DOCX → text + extract ảnh nhúng (qua python-docx).

    Full mode: trả `[{"page": int, "text": str, "images": list}]` — schema
    đồng nhất với parse_pdf. Khi `doc_id` rỗng (vd preview), bỏ qua image
    extraction.

    Preview mode (`max_paragraphs > 0`): trả `str` plain text cho metadata gen,
    KHÔNG extract ảnh — backward compat với caller hiện tại.

    DOCX không có page boundary thật → ảnh sẽ gom vào page=1; caption dùng
    context của page=1 (text Docling parse cho phần đầu doc).
    """
    # Preview path: dùng python-docx trực tiếp, break sớm — bypass Docling
    # để tiết kiệm thời gian (Docling convert toàn bộ docx, kể cả khi chỉ cần
    # vài đoạn đầu cho metadata).
    if max_paragraphs and max_paragraphs > 0:
        try:
            from docx import Document
            doc = Document(str(path))
            parts: list[str] = []
            for para in doc.paragraphs:
                if len(parts) >= max_paragraphs:
                    break
                text = para.text.strip()
                if text:
                    parts.append(text)
            return "\n\n".join(parts)
        except Exception:
            logger.exception("python-docx preview failed for %s", path.name)
            return ""

    def _attach_docx_images(pages: list[dict]) -> list[dict]:
        if not doc_id:
            for p in pages:
                p.setdefault("images", [])
            return pages
        page_texts = {p["page"]: p["text"] for p in pages}
        images_by_page = _extract_caption_upload_images_docx(path, doc_id, page_texts)
        for p in pages:
            p["images"] = images_by_page.get(p["page"], [])
        return pages

    # Tier 1: Docling — per-page list (split theo heading hoặc paragraph block)
    pages = _parse_with_docling(path)
    if pages:
        return _attach_docx_images(pages)

    # Tier 2: python-docx fallback — chỉ trả 1 page với toàn bộ text + heading.
    try:
        from docx import Document
        doc = Document(str(path))
        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name if para.style else ""
            if style.startswith("Heading"):
                level = "".join(filter(str.isdigit, style)) or "1"
                parts.append("#" * int(level) + " " + text)
            else:
                parts.append(text)
        full_text = "\n\n".join(parts)
        if not full_text.strip():
            return [{"page": 1, "text": "", "images": []}]
        return _attach_docx_images([{"page": 1, "text": full_text}])
    except Exception:
        logger.exception("python-docx failed for %s", path.name)
        return [{"page": 1, "text": "", "images": []}]


def parse_pptx(path: Path, max_slides: int | None = None) -> list[dict] | str:
    # Preview path: chỉ N slide đầu, dùng python-pptx trực tiếp — bypass Docling
    if max_slides and max_slides > 0:
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            result: list[dict] = []
            for i, slide in enumerate(prs.slides, 1):
                if i > max_slides:
                    break
                texts = [
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                if texts:
                    result.append({"page": i, "text": "\n".join(texts)})
            return result or ""
        except Exception:
            logger.exception("python-pptx preview failed for %s", path.name)
            return ""
    pages = _parse_with_docling(path)
    if pages:
        return pages
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        result: list[dict] = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                result.append({"page": i, "text": "\n".join(texts)})
        return result or ""
    except Exception:
        logger.exception("python-pptx failed for %s", path.name)
        return ""


def parse_txt(path: Path, max_bytes: int | None = None) -> str:
    if max_bytes and max_bytes > 0:
        try:
            with path.open("rb") as f:
                raw = f.read(max_bytes)
            for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
                try:
                    return raw.decode(encoding)
                except UnicodeDecodeError:
                    continue
            return raw.decode("utf-8", errors="replace")
        except Exception:
            logger.exception("parse_txt preview failed for %s", path.name)
            return ""
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except (UnicodeDecodeError, ValueError):
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def parse_md(path: Path, max_bytes: int | None = None) -> str:
    return parse_txt(path, max_bytes=max_bytes)


# --- Main entry ---

# Heuristics cho preview: 1 "page" tương đương ~5 paragraph DOCX, ~50 row XLSX,
# ~10KB TXT/MD. Mục đích: parse đủ context cho Haiku gen metadata, không hơn.
_PREVIEW_DOCX_PARAS_PER_PAGE = 5
_PREVIEW_XLSX_ROWS_PER_PAGE = 50
_PREVIEW_TXT_BYTES_PER_PAGE = 10240


def parse(path: Union[str, Path], max_pages: int | None = None) -> dict:
    """Parse file → dict {doc_id, source, content, mime, uploaded_at}.

    Khi `max_pages` set, mỗi parser chỉ đọc đoạn đầu — dùng cho preview metadata.
    Khi None (mặc định), parse đầy đủ — dùng cho ingest thật vào Qdrant.
    """
    path = Path(path)
    if not path.exists():
        logger.error("File not found: %s", path)
        return {"doc_id": "0" * 16, "source": path.name, "content": "",
                "mime": "application/octet-stream", "uploaded_at": datetime.now(timezone.utc).isoformat()}
    if path.stat().st_size == 0:
        logger.warning("Empty file: %s", path.name)
        return {"doc_id": _sha256(path), "source": path.name, "content": "",
                "mime": "application/octet-stream", "uploaded_at": datetime.now(timezone.utc).isoformat()}

    suffix = path.suffix.lower()
    doc_id = _sha256(path)
    uploaded_at = datetime.now(timezone.utc).isoformat()

    if suffix == ".pdf":
        content = parse_pdf(path, doc_id=doc_id, max_pages=max_pages)
        mime = "application/pdf"
    elif suffix in (".docx", ".doc"):
        content = parse_docx(
            path,
            doc_id=doc_id,
            max_paragraphs=(max_pages * _PREVIEW_DOCX_PARAS_PER_PAGE) if max_pages else None,
        )
        mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    elif suffix in (".pptx", ".ppt"):
        content = parse_pptx(path, max_slides=max_pages)
        mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    elif suffix == ".xlsx":
        content = parse_xlsx(
            path,
            doc_id=doc_id,
            max_rows=(max_pages * _PREVIEW_XLSX_ROWS_PER_PAGE) if max_pages else None,
        )
        mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    elif suffix == ".md":
        content = parse_md(
            path,
            max_bytes=(max_pages * _PREVIEW_TXT_BYTES_PER_PAGE) if max_pages else None,
        )
        mime = "text/markdown"
    else:
        content = parse_txt(
            path,
            max_bytes=(max_pages * _PREVIEW_TXT_BYTES_PER_PAGE) if max_pages else None,
        )
        mime = mimetypes.guess_type(path.name)[0] or "text/plain"

    return {"doc_id": doc_id, "source": path.name, "content": content,
            "mime": mime, "uploaded_at": uploaded_at}
